#!/usr/bin/env python3
"""
extract_pptx_structure.py

Extract the structure of an existing .pptx file into a JSON description.
Used for the "read → understand → rebuild in PptxGenJS" workflow.

Usage:
    python3 extract_pptx_structure.py input.pptx                     # prints JSON to stdout
    python3 extract_pptx_structure.py input.pptx -o structure.json   # saves to file
    python3 extract_pptx_structure.py input.pptx --extract-images    # also exports images

Dependencies:
    pip install python-pptx Pillow
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def emu_to_inches(emu_val):
    """Convert EMU to inches."""
    if emu_val is None:
        return None
    return round(emu_val / 914400, 4)


def extract_text_properties(text_frame):
    """Extract text content and formatting from a text frame."""
    paragraphs = []
    for para in text_frame.paragraphs:
        runs = []
        for run in para.runs:
            run_info = {
                "text": run.text,
            }
            font = run.font
            if font.size:
                run_info["fontSize"] = round(font.size.pt, 1)
            if font.bold:
                run_info["bold"] = True
            if font.italic:
                run_info["italic"] = True
            if font.color and font.color.rgb:
                run_info["color"] = str(font.color.rgb)
            if font.name:
                run_info["fontFace"] = font.name
            runs.append(run_info)

        para_info = {
            "text": para.text,
            "runs": runs,
        }
        if para.alignment is not None:
            para_info["alignment"] = str(para.alignment)
        paragraphs.append(para_info)

    return {
        "fullText": text_frame.text,
        "paragraphs": paragraphs,
    }


def extract_shape(shape, slide_index, shape_index, extract_images, output_dir):
    """Extract information from a single shape."""
    info = {
        "index": shape_index,
        "name": shape.name,
        "type": str(shape.shape_type) if hasattr(shape, "shape_type") else "unknown",
        "position": {
            "x": emu_to_inches(shape.left),
            "y": emu_to_inches(shape.top),
            "w": emu_to_inches(shape.width),
            "h": emu_to_inches(shape.height),
        },
    }

    # Rotation
    if hasattr(shape, "rotation") and shape.rotation:
        info["rotation"] = shape.rotation

    # Text content
    if shape.has_text_frame:
        info["text"] = extract_text_properties(shape.text_frame)

    # Image
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["type"] = "image"
        image = shape.image
        info["image"] = {
            "content_type": image.content_type,
            "width": image.size[0] if hasattr(image, "size") else None,
            "height": image.size[1] if hasattr(image, "size") else None,
        }

        if extract_images and output_dir:
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            img_filename = f"slide{slide_index + 1}_shape{shape_index}.{ext}"
            img_path = os.path.join(output_dir, img_filename)
            with open(img_path, "wb") as f:
                f.write(image.blob)
            info["image"]["extracted_path"] = img_path

    # Table
    if shape.has_table:
        info["type"] = "table"
        table = shape.table
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text)
            rows.append(cells)
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "data": rows,
        }

    # Chart
    if shape.has_chart:
        info["type"] = "chart"
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
        }

    # Group shape
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["type"] = "group"
        info["children"] = []
        for i, child in enumerate(shape.shapes):
            child_info = extract_shape(child, slide_index, i, extract_images, output_dir)
            info["children"].append(child_info)

    # Fill
    if hasattr(shape, "fill"):
        fill = shape.fill
        if fill.type is not None:
            info["fill"] = {"type": str(fill.type)}

    return info


def extract_slide(slide, slide_index, extract_images, output_dir):
    """Extract information from a single slide."""
    slide_info = {
        "index": slide_index,
        "layout": slide.slide_layout.name if slide.slide_layout else None,
        "elements": [],
    }

    # Background
    if slide.background and slide.background.fill:
        bg_fill = slide.background.fill
        if bg_fill.type is not None:
            slide_info["background"] = {"type": str(bg_fill.type)}

    # Shapes
    for i, shape in enumerate(slide.shapes):
        shape_info = extract_shape(shape, slide_index, i, extract_images, output_dir)
        slide_info["elements"].append(shape_info)

    # Notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        slide_info["notes"] = slide.notes_slide.notes_text_frame.text.strip()

    return slide_info


def extract_pptx(pptx_path, extract_images=False, output_dir=None):
    """Extract the full structure of a .pptx file."""
    prs = Presentation(pptx_path)

    if extract_images and output_dir:
        os.makedirs(output_dir, exist_ok=True)

    structure = {
        "file": os.path.basename(pptx_path),
        "slide_width": emu_to_inches(prs.slide_width),
        "slide_height": emu_to_inches(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    # Slide masters and layouts
    layouts = []
    for layout in prs.slide_layouts:
        layouts.append({
            "name": layout.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name}
                for ph in layout.placeholders
            ],
        })
    structure["available_layouts"] = layouts

    # Slides
    for i, slide in enumerate(prs.slides):
        slide_info = extract_slide(slide, i, extract_images, output_dir)
        structure["slides"].append(slide_info)

    return structure


def main():
    parser = argparse.ArgumentParser(description="Extract .pptx structure to JSON")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("-o", "--output", help="Output JSON file (default: stdout)")
    parser.add_argument("--extract-images", action="store_true",
                        help="Extract images to a directory")
    parser.add_argument("--image-dir", default="extracted_assets",
                        help="Directory for extracted images (default: extracted_assets)")
    parser.add_argument("--pretty", action="store_true", default=True,
                        help="Pretty-print JSON output (default: True)")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: File not found: {args.input}", file=sys.stderr)
        sys.exit(1)

    structure = extract_pptx(
        args.input,
        extract_images=args.extract_images,
        output_dir=args.image_dir if args.extract_images else None,
    )

    indent = 2 if args.pretty else None
    json_str = json.dumps(structure, indent=indent, ensure_ascii=False)

    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(json_str)
        print(f"✅ Structure extracted to {args.output}", file=sys.stderr)
        if args.extract_images:
            img_count = sum(
                1 for s in structure["slides"]
                for e in s["elements"]
                if e.get("type") == "image" and "extracted_path" in e.get("image", {})
            )
            print(f"✅ Extracted {img_count} images to {args.image_dir}/", file=sys.stderr)
    else:
        print(json_str)


if __name__ == "__main__":
    main()
